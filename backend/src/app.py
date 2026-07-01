from fastapi import FastAPI, UploadFile, File
from fastapi.middleware.cors import CORSMiddleware
import google.generativeai as genai
from pypdf import PdfReader
from pptx import Presentation
from docx import Document
from io import BytesIO
import uvicorn
from dotenv import load_dotenv
load_dotenv()
import os
app = FastAPI()

app.add_middleware(
    CORSMiddleware,
    allow_origins=[
        "https://ai-summarizer-ecru.vercel.app",
        "http://localhost:5173", 
    ],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)
genai.configure(api_key=os.getenv("GEMINI_API_KEY"))
model = genai.GenerativeModel("gemini-2.5-flash")
@app.post("/summarize")
async def summarizing(file: UploadFile = File(...)):
    text = ""

    if file.filename.endswith(".pdf"):
        content = await file.read()
        reader= PdfReader(BytesIO(content))
        for page in reader.pages:
            text+=page.extract_text() or ""

    elif file.filename.endswith(".docx"):
        content = await file.read()
        reader= Document(BytesIO(content))
        for para in reader.paragraphs:
            text+=para.text + "\n"

    elif file.filename.endswith(".pptx"):
        content = await file.read()
        presentation = Presentation(BytesIO(content))
        for slide in presentation.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
                    
    prompt = f"""
    You are an expert study assistant.

    Analyze the following notes and create a concise but comprehensive summary.

    Instructions:
    - Include all important concepts, definitions, formulas, and facts.
    - Remove repetition and unnecessary details.
    - Use clear headings and bullet points.
    - Keep the summary brief and easy to revise from.
    - If the notes contain examples, mention them only if they help explain a concept.
    - If there are formulas or code snippets, preserve them exactly.

    Output format:

    # Short Summary
    A 4-6 sentence overview of the entire document.

    # Key Concepts
    - Concept 1: Explanation
    - Concept 2: Explanation
    - ...

    # Important Points
    - Point 1
    - Point 2
    - ...

    # Formulas / Definitions (if any)
    - ...

    # Potential Exam Questions
    1. ...
    2. ...
    3. ...

    Notes:
    {text}
    """
    response = model.generate_content(prompt)
    print(response.text)

    return {
        "summ": response.text
    }



if __name__ == "__main__":

    uvicorn.run(
        app,
        host="0.0.0.0",
        port=8000
    )

